"""Generates one test document per supported format into fixtures/out/,
all with accented + spaced filenames, plus a corrupted PDF and a ZIP.

Run with the project venv: python\\.venv\\Scripts\\python.exe fixtures\\make_fixtures.py
Dev-only deps (requirements-dev.txt): python-docx, fpdf2, xlwt.
"""

import json
import os
import zipfile

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'out')
os.makedirs(OUT, exist_ok=True)


def path(name):
    return os.path.join(OUT, name)


def make_docx():
    import docx
    document = docx.Document()
    document.add_heading('Relatório Anual', level=1)
    document.add_paragraph('Este é um parágrafo com acentuação: ção, ã, é, í, õ, ü.')
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = 'Região'
    table.rows[0].cells[1].text = 'Receita'
    table.rows[1].cells[0].text = 'São Paulo'
    table.rows[1].cells[1].text = 'R$ 1.234,56'
    document.save(path('relatório anual (versão 2).docx'))


def make_pptx():
    from pptx import Presentation
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = 'Apresentação de Vendas'
    slide.placeholders[1].text = 'Trimestre três — resultados preliminares'
    bullet_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    bullet_slide.shapes.title.text = 'Metas'
    bullet_slide.placeholders[1].text = 'Crescimento de 10%\nExpansão regional\nRetenção de clientes'
    presentation.save(path('apresentação de vendas.pptx'))


def make_xlsx():
    import openpyxl
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = 'Orçamento'
    sheet.append(['Item', 'Descrição', 'Valor'])
    sheet.append(['A1', 'Manutenção prédio', 1500.5])
    sheet.append(['B2', 'Refeições da equipe', 320.0])
    workbook.save(path('planilha de orçamento.xlsx'))


def make_xls():
    import xlwt
    workbook = xlwt.Workbook(encoding='utf-8')
    sheet = workbook.add_sheet('Antiga')
    for column, header in enumerate(['Código', 'Município', 'População']):
        sheet.write(0, column, header)
    sheet.write(1, 0, 1)
    sheet.write(1, 1, 'Brasília')
    sheet.write(1, 2, 2817068)
    workbook.save(path('tabela antiga.xls'))


def make_pdf():
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font('helvetica', size=14)
    pdf.cell(text='Documento de Teste em PDF')
    pdf.ln(12)
    pdf.set_font('helvetica', size=11)
    pdf.multi_cell(w=180, text='Conteúdo com acentuação: ção, ã, é. Segunda linha do parágrafo.')
    pdf.output(path('documento em pdf.pdf'))


def make_texty():
    with open(path('dados de vendas.csv'), 'w', encoding='utf-8', newline='') as handle:
        handle.write('região,produto,valor\nSão Paulo,Cadeira,199.90\nBrasília,Mesa,349.00\n')
    with open(path('configuração.json'), 'w', encoding='utf-8') as handle:
        json.dump({'aplicação': 'teste', 'versão': 2, 'ativo': True}, handle, ensure_ascii=False, indent=2)
    with open(path('notas fiscais.xml'), 'w', encoding='utf-8') as handle:
        handle.write('<?xml version="1.0" encoding="UTF-8"?>\n<notas><nota id="1"><cliente>João</cliente><valor>10.50</valor></nota></notas>\n')
    with open(path('página inicial.html'), 'w', encoding='utf-8') as handle:
        handle.write('<!doctype html><html lang="pt"><head><meta charset="utf-8"><title>Página de Teste</title></head>'
                     '<body><h1>Título Principal</h1><p>Um parágrafo com <strong>negrito</strong> e acentuação.</p>'
                     '<ul><li>Item um</li><li>Item dois</li></ul></body></html>\n')
    with open(path('leia-me primeiro.txt'), 'w', encoding='utf-8') as handle:
        handle.write('Arquivo de texto simples.\nSegunda linha com acentuação: ção.\n')


def make_zip():
    with zipfile.ZipFile(path('arquivos compactados.zip'), 'w') as archive:
        archive.write(path('leia-me primeiro.txt'), 'leia-me primeiro.txt')
        archive.write(path('dados de vendas.csv'), 'dados de vendas.csv')
        archive.write(path('relatório anual (versão 2).docx'), 'relatório anual (versão 2).docx')


def make_corrupted():
    with open(path('corrompido.pdf'), 'wb') as handle:
        handle.write(b'%PDF-1.7\n' + os.urandom(2048))


if __name__ == '__main__':
    make_docx()
    make_pptx()
    make_xlsx()
    make_xls()
    make_pdf()
    make_texty()
    make_zip()
    make_corrupted()
    for name in sorted(os.listdir(OUT)):
        print(name, os.path.getsize(path(name)))
